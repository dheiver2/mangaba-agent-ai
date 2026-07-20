"""Leitura de documentos Office e CSV (complemento do read_pdf)."""

import csv as csv_mod

from app.tool.base import BaseTool, ToolResult
from app.tool.gateway_media import _resolve

_PREVIEW = 5000
_MAX_ROWS = 50


class ReadOffice(BaseTool):
    name: str = "read_office"
    description: str = (
        "Extrai o conteúdo de documentos .docx (texto e tabelas), .xlsx (planilhas, "
        "primeiras linhas de cada aba), .pptx (texto dos slides) e .csv (primeiras "
        "linhas). Para PDF use read_pdf. Caminhos relativos caem no workspace."
    )
    parameters: dict = {
        "type": "object",
        "properties": {
            "file_path": {
                "type": "string",
                "description": "Caminho do documento (.docx, .xlsx, .pptx ou .csv)",
            },
        },
        "required": ["file_path"],
    }

    async def execute(self, file_path: str, **kwargs) -> ToolResult:
        path = _resolve(file_path)
        if not path.is_file():
            return ToolResult(error=f"Arquivo não encontrado: {path}")
        ext = path.suffix.lower()

        try:
            if ext == ".docx":
                from docx import Document

                doc = Document(str(path))
                parts = [p.text for p in doc.paragraphs if p.text.strip()]
                for ti, table in enumerate(doc.tables, 1):
                    rows = ["\t".join(c.text for c in row.cells) for row in table.rows]
                    parts.append(f"[Tabela {ti}]\n" + "\n".join(rows))
                text = "\n".join(parts)

            elif ext == ".xlsx":
                from openpyxl import load_workbook

                wb = load_workbook(str(path), read_only=True, data_only=True)
                parts = []
                for ws in wb.worksheets:
                    lines = []
                    for ri, row in enumerate(ws.iter_rows(values_only=True)):
                        if ri >= _MAX_ROWS:
                            lines.append(f"… (+{ws.max_row - _MAX_ROWS} linhas)")
                            break
                        lines.append(
                            "\t".join("" if c is None else str(c) for c in row)
                        )
                    parts.append(
                        f"=== Aba '{ws.title}' ({ws.max_row}x{ws.max_column}) ===\n"
                        + "\n".join(lines)
                    )
                wb.close()
                text = "\n\n".join(parts)

            elif ext == ".pptx":
                from pptx import Presentation

                prs = Presentation(str(path))
                parts = []
                for si, slide in enumerate(prs.slides, 1):
                    texts = [
                        sh.text.strip()
                        for sh in slide.shapes
                        if getattr(sh, "has_text_frame", False) and sh.text.strip()
                    ]
                    parts.append(f"--- Slide {si} ---\n" + "\n".join(texts))
                text = "\n".join(parts)

            elif ext in (".csv", ".tsv"):
                delim = "\t" if ext == ".tsv" else ","
                with open(path, newline="", encoding="utf-8", errors="replace") as fh:
                    reader = csv_mod.reader(fh, delimiter=delim)
                    rows = []
                    total = 0
                    for row in reader:
                        total += 1
                        if total <= _MAX_ROWS:
                            rows.append("\t".join(row))
                text = "\n".join(rows)
                if total > _MAX_ROWS:
                    text += f"\n… (+{total - _MAX_ROWS} linhas; total {total})"

            else:
                return ToolResult(
                    error=f"Extensão não suportada: {ext} (aceito: .docx .xlsx .pptx .csv .tsv)"
                )
        except Exception as e:
            return ToolResult(error=f"Falha ao ler {path.name}: {e}")

        if not text.strip():
            return ToolResult(output=f"{path.name} não contém texto extraível")
        preview = text[:_PREVIEW] + ("… (truncado)" if len(text) > _PREVIEW else "")
        return ToolResult(output=f"Conteúdo de {path.name}:\n{preview}")
